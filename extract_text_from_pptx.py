#!/usr/bin/python3
import sys
import zipfile
import olefile
from io import BytesIO
from pptx import Presentation
from pptx.exc import PackageNotFoundError

def repair_pptx_file(pptx_file):
    with zipfile.ZipFile(pptx_file, 'r') as original:
        repaired_pptx_data = BytesIO()
        with zipfile.ZipFile(repaired_pptx_data, 'w') as repaired:
            for item in original.infolist():
                try:
                    data = original.read(item.filename)
                    repaired.writestr(item, data)
                except zipfile.BadZipFile:
                    print(f"Warning: Skipping bad file part '{item.filename}' in '{pptx_file}'")
        return BytesIO(repaired_pptx_data.getvalue())

def extract_text_from_presentation(pptx_file):
    try:
        ppt = Presentation(pptx_file)
    except zipfile.BadZipFile:
        print(f"Warning: Encountered a bad zipfile in '{pptx_file}'. Attempting to repair...")
        repaired_pptx_file = repair_pptx_file(pptx_file)
        ppt = Presentation(repaired_pptx_file)

    text_list = []

    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text_list.append(paragraph.text)

    return '\n'.join(text_list)

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print('Usage: python extract_text_from_pptx.py <pptx_file_path>')
        sys.exit(1)

    pptx_file = sys.argv[1]

    try:
        text = extract_text_from_presentation(pptx_file)
        print(text)
    except PackageNotFoundError:
        print(f"Error: Unable to open '{pptx_file}'. Make sure the file exists and is a valid PPTX file.")
